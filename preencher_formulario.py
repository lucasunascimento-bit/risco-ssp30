# ============================================================
# preencher_formulario.py
# Lê os 41 PPTXs de Analise Driver e gera um relatório HTML
# completo com todos os drivers para impressão/PDF.
#
# Uso: python preencher_formulario.py
# Dep: pip install python-pptx
# ============================================================

import os, re, webbrowser, tempfile
from datetime import datetime
from pptx import Presentation

PASTA_PPTX = r'C:\Users\lucasn\Downloads\Analise Driver'
RESPONSAVEL = 'Lucas de Oliveira Nascimento'
OPERACAO    = 'SSP30'

# ============================================================
# EXTRAÇÃO
# ============================================================

def _textos_slide(slide):
    """Retorna lista de strings não-vazias de um slide."""
    out = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    out.append(t)
    return out

def extrair_dados(pptx_path):
    prs = Presentation(pptx_path)
    s1  = '\n'.join(_textos_slide(prs.slides[0]))
    s2  = '\n'.join(_textos_slide(prs.slides[1])) if len(prs.slides) > 1 else ''
    s3  = '\n'.join(_textos_slide(prs.slides[2])) if len(prs.slides) > 2 else ''
    full = s1 + '\n' + s2

    dados = {
        'r_id':         '',
        'responsavel':  RESPONSAVEL,
        'r_placa':      '',
        'r_mlp':        '',
        'r_valor':      '',
        'data_reporte': '',
        'ocorrencias':  [],
        '_raw_s2':      s2,
    }

    # Driver ID
    m = re.search(r'Driver\s*ID:\s*(\d+)', full)
    if m: dados['r_id'] = m.group(1)

    # Placa (formato BR: ABC1234 ou ABC1A23)
    m = re.search(r'Placa:\s*([A-Z]{3}\d[A-Z0-9]\d{2})', full)
    if m: dados['r_placa'] = m.group(1).strip()

    # MLP (transportadora) — vem no slide 2 como "MLP: NOME   |"
    m = re.search(r'MLP:\s*([^|]+?)(?:\s+\||\s*$)', s2, re.MULTILINE)
    if m: dados['r_mlp'] = m.group(1).strip()

    # Total BPP
    m = re.search(r'Total\s*BPP:\s*\$?([\d\.,]+)', s2)
    if m: dados['r_valor'] = m.group(1).replace(',', '.').strip()

    # Data solicitação
    m = re.search(r'Data\s*solicitacao:\s*([\d/\-]+)', s1)
    if m: dados['data_reporte'] = m.group(1)

    # Ocorrências (tabela no slide 2)
    linhas = [l for l in s2.splitlines() if l.strip()]
    cabecalho_idx = None
    for i, l in enumerate(linhas):
        if 'ID SHIPMENT' in l.upper() or 'SHIPMENT' in l.upper():
            cabecalho_idx = i
            break
    if cabecalho_idx is not None:
        for l in linhas[cabecalho_idx + 1:]:
            partes = l.split()
            if partes and re.match(r'\d{10,}', partes[0] if len(partes) > 1 else ''):
                dados['ocorrencias'].append(l.strip())
            elif re.search(r'\d{10,}', l):
                dados['ocorrencias'].append(l.strip())

    return dados

# ============================================================
# HTML REPORT
# ============================================================

_CSS = """
body{font-family:Arial,sans-serif;background:#f5f5f5;margin:0;padding:20px;color:#1a1a1a}
h1{color:#2563eb;margin-bottom:4px}
.meta{color:#6b7280;font-size:13px;margin-bottom:24px}
.driver-card{background:#fff;border:1px solid #e5e7eb;border-radius:8px;margin-bottom:20px;
             page-break-inside:avoid;box-shadow:0 1px 3px rgba(0,0,0,.08)}
.card-header{background:#1e3a8a;color:#fff;padding:12px 16px;border-radius:8px 8px 0 0;
             display:flex;justify-content:space-between;align-items:center}
.card-header h2{margin:0;font-size:15px}
.driver-id{font-family:monospace;font-size:13px;opacity:.8}
.card-body{padding:16px}
.fields{display:grid;grid-template-columns:repeat(3,1fr);gap:12px;margin-bottom:12px}
.field label{font-size:11px;color:#6b7280;display:block;margin-bottom:2px;text-transform:uppercase;letter-spacing:.5px}
.field span{font-weight:600;font-size:14px}
.ocorr-table{width:100%;border-collapse:collapse;font-size:12px;margin-top:8px}
.ocorr-table th{background:#f3f4f6;padding:6px 8px;text-align:left;font-size:11px;
                color:#6b7280;border-bottom:2px solid #e5e7eb;text-transform:uppercase}
.ocorr-table td{padding:6px 8px;border-bottom:1px solid #f3f4f6}
.ocorr-table tr:last-child td{border-bottom:none}
.summary-bar{background:#fff;border:1px solid #e5e7eb;border-radius:8px;padding:16px;
             margin-bottom:24px;display:flex;gap:32px;align-items:center}
.stat{text-align:center}
.stat .num{font-size:28px;font-weight:700;color:#2563eb}
.stat .lbl{font-size:11px;color:#6b7280;text-transform:uppercase}
.erro{color:#ef4444;font-size:12px;padding:8px;background:#fef2f2;border-radius:4px}
@media print{.driver-card{box-shadow:none;border:1px solid #ccc}}
"""

def _card_ocorrencias(ocorrs):
    if not ocorrs:
        return '<p style="color:#9ca3af;font-size:12px;margin:0">Nenhuma ocorrência extraída</p>'
    rows = ''.join(f'<tr><td colspan="7">{o}</td></tr>' for o in ocorrs)
    return f'''<table class="ocorr-table">
        <thead><tr>
            <th>Linha raw (extraída do PPTX)</th>
        </tr></thead>
        <tbody>{rows}</tbody>
    </table>'''

def gerar_html(registros):
    agora   = datetime.now().strftime('%d/%m/%Y %H:%M')
    n_ok    = sum(1 for r in registros if not r.get('erro'))
    n_erro  = sum(1 for r in registros if r.get('erro'))
    gmv_total = 0.0
    for r in registros:
        d = r.get('dados', {})
        try:
            gmv_total += float(d.get('r_valor', 0) or 0)
        except Exception:
            pass

    cards = ''
    for i, reg in enumerate(registros, 1):
        nome = reg['nome']
        if reg.get('erro'):
            cards += f'''<div class="driver-card">
                <div class="card-header" style="background:#7f1d1d">
                    <h2>{i}. {nome}</h2>
                </div>
                <div class="card-body"><p class="erro">Erro ao extrair: {reg["erro"]}</p></div>
            </div>'''
            continue
        d = reg['dados']
        ocorr_html = _card_ocorrencias(d.get('ocorrencias', []))
        cards += f'''<div class="driver-card">
            <div class="card-header">
                <h2>{i}. {d.get("r_mlp") or "—"}</h2>
                <span class="driver-id">ID {d.get("r_id","?")}</span>
            </div>
            <div class="card-body">
                <div class="fields">
                    <div class="field"><label>Driver ID</label><span>{d.get("r_id","—")}</span></div>
                    <div class="field"><label>Placa</label><span>{d.get("r_placa","—")}</span></div>
                    <div class="field"><label>MLP / Transportadora</label><span>{d.get("r_mlp","—")}</span></div>
                    <div class="field"><label>Total BPP</label><span style="color:#ef4444">${d.get("r_valor","—")}</span></div>
                    <div class="field"><label>Data Solicitação</label><span>{d.get("data_reporte","—")}</span></div>
                    <div class="field"><label>Responsável</label><span style="font-size:12px">{d.get("responsavel","—")}</span></div>
                </div>
                {ocorr_html}
            </div>
        </div>'''

    return f'''<!DOCTYPE html>
<html lang="pt-BR">
<head>
<meta charset="utf-8">
<title>Relatório de Investigação — {OPERACAO}</title>
<style>{_CSS}</style>
</head>
<body>
<h1>Relatório de Investigação LP — {OPERACAO}</h1>
<p class="meta">Gerado em {agora} · {len(registros)} drivers · {n_ok} ok · {n_erro} com erro</p>
<div class="summary-bar">
    <div class="stat"><div class="num">{len(registros)}</div><div class="lbl">Drivers</div></div>
    <div class="stat"><div class="num">{n_ok}</div><div class="lbl">Extraídos OK</div></div>
    <div class="stat"><div class="num" style="color:#ef4444">${gmv_total:,.2f}</div><div class="lbl">Total BPP</div></div>
    <div class="stat"><div class="num">{n_erro}</div><div class="lbl">Erros extração</div></div>
</div>
{cards}
</body>
</html>'''

# ============================================================
# MAIN
# ============================================================

def main():
    arquivos = sorted([
        f for f in os.listdir(PASTA_PPTX) if f.lower().endswith('.pptx')
    ])
    print(f'Encontrados {len(arquivos)} PPTXs em {PASTA_PPTX}')

    registros = []
    for nome in arquivos:
        caminho = os.path.join(PASTA_PPTX, nome)
        try:
            dados = extrair_dados(caminho)
            registros.append({'nome': nome, 'dados': dados})
            print(f'  OK  Driver {dados["r_id"] or "?"} | {dados["r_mlp"] or "?"} | ${dados["r_valor"] or "?"}')
        except Exception as e:
            registros.append({'nome': nome, 'erro': str(e)})
            print(f'  ERR {nome}: {e}')

    html = gerar_html(registros)
    saida = os.path.join(PASTA_PPTX, f'relatorio_investigacao_{datetime.now().strftime("%Y%m%d_%H%M")}.html')
    with open(saida, 'w', encoding='utf-8') as f:
        f.write(html)
    print(f'\nRelatório salvo: {saida}')
    webbrowser.open('file:///' + saida.replace('\\', '/'))

if __name__ == '__main__':
    main()
