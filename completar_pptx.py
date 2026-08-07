"""
completar_pptx.py — Adiciona slides de continuação nos PPTXs de Analise Driver
para drivers que têm mais de 32 SHPs no CSV enriquecido.

Uso: python completar_pptx.py
Saída: copia modificada em Downloads\Analise Driver\Completo\
"""

import copy, glob, re, os
from pathlib import Path
from datetime import datetime
import pandas as pd
from pptx import Presentation
from pptx.util import Pt

# ── Config ────────────────────────────────────────────────────────────────────
PASTA_PPTX  = Path(r'C:\Users\lucasn\Downloads\Analise Driver')
PASTA_OUT   = PASTA_PPTX / 'Completo'
ROWS_PER_SLIDE = 32

MESES_PT = {
    1:'Janeiro', 2:'Fevereiro', 3:'Março', 4:'Abril', 5:'Maio', 6:'Junho',
    7:'Julho', 8:'Agosto', 9:'Setembro', 10:'Outubro', 11:'Novembro', 12:'Dezembro',
}

# ── Carregar CSV ───────────────────────────────────────────────────────────────
def carregar_csv():
    csvs = glob.glob(str(Path(r'C:\Users\lucasn\Downloads') / 'KPI*com_driver*.csv'))
    if not csvs:
        raise FileNotFoundError('CSV *com_driver* não encontrado em Downloads')
    csv = sorted(csvs, key=os.path.getmtime, reverse=True)[0]
    print(f'CSV: {Path(csv).name}')
    df = pd.read_csv(csv, dtype={'SHIPMENT_ID': str, 'DRIVER_ID': str})
    df['SHIPMENT_ID'] = df['SHIPMENT_ID'].str.strip()
    df['DRIVER_ID']   = df['DRIVER_ID'].str.strip()
    before = len(df)
    df = df.drop_duplicates(subset=['SHIPMENT_ID'], keep='first')
    after = len(df)
    if before != after:
        print(f'  Dedup: {before - after} SHIPMENT_IDs duplicados removidos')
    return df

# ── Helpers de formatação ─────────────────────────────────────────────────────
def _mes(data_str):
    """'31/07/2026' → 'Julho'"""
    try:
        d = datetime.strptime(data_str.strip(), '%d/%m/%Y')
        return MESES_PT.get(d.month, data_str)
    except Exception:
        return data_str or ''

def _bpp(v):
    """float → '$94,76'"""
    try:
        return f'${float(v):,.2f}'.replace('.', 'X').replace(',', '.').replace('X', ',')
    except Exception:
        return str(v)

def _row_vals(r):
    """Retorna lista de 7 valores para preencher um row do PPTX."""
    return [
        _mes(str(r.get('DATA_BQ', '') or '')),              # col 0 MES
        str(r.get('SHIPMENT_ID', '') or ''),                # col 1 ID SHIPMENT
        str(r.get('DATA_BQ', '') or ''),                    # col 2 DATA BPP
        str(r.get('CLASSIFICACAO_BQ', '') or ''),           # col 3 CLASSIFICACAO LG
        str(r.get('TIPO_DAMAGED_LG', '') or ''),            # col 4 DANO
        str(r.get('DOM_DOMAIN_ID', '') or ''),              # col 5 DOMINIO
        _bpp(r.get('BPP', 0)),                              # col 6 BPP CASHOUT
    ]

# ── Copiar slide ──────────────────────────────────────────────────────────────
def copy_slide(prs, source_idx):
    """Duplica um slide e adiciona ao final da apresentação."""
    src = prs.slides[source_idx]
    new_slide = prs.slides.add_slide(src.slide_layout)
    sp_tree = new_slide.shapes._spTree
    for child in list(sp_tree):
        sp_tree.remove(child)
    for child in src.shapes._spTree:
        sp_tree.append(copy.deepcopy(child))
    return new_slide


def move_slide(prs, from_idx, to_idx):
    """Move slide de from_idx para to_idx no XML da apresentação."""
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    slide_elem = slides[from_idx]
    xml_slides.remove(slide_elem)
    xml_slides.insert(to_idx, slide_elem)

# ── Preencher um slide com uma janela de dados ────────────────────────────────
def preencher_slide(slide, rows_batch, slide_num, total_slides):
    """
    rows_batch: lista de dicts (linhas do CSV)
    Atualiza shapes c1007-c1230 com os dados.
    """
    shapes = {s.name: s for s in slide.shapes if s.has_text_frame}

    # Atualizar título/InfoBar para indicar continuação
    for nome in ('Title', 'InfoBar'):
        if nome in shapes:
            tf = shapes[nome].text_frame
            for para in tf.paragraphs:
                for run in para.runs:
                    if nome == 'Title' and slide_num > 1:
                        # Acrescentar indicador de página ao título
                        if f'(pág. {slide_num}' not in run.text:
                            run.text = run.text.rstrip() + f' (pág. {slide_num}/{total_slides})'
                    break

    for row_idx in range(ROWS_PER_SLIDE):
        base = 1007 + row_idx * 7
        if row_idx < len(rows_batch):
            vals = _row_vals(rows_batch[row_idx])
        else:
            vals = [''] * 7  # limpar células vazias

        for col in range(7):
            nome = f'c{base + col}'
            if nome in shapes:
                tf = shapes[nome].text_frame
                # Preservar formatação: atualiza só o run text
                for para in tf.paragraphs:
                    for run in para.runs:
                        run.text = vals[col]
                        break
                    else:
                        # Sem run: criar
                        from pptx.util import Pt
                        run = para.add_run()
                        run.text = vals[col]

# ── Processar um PPTX ─────────────────────────────────────────────────────────
def processar(pptx_path, df_driver, out_dir):
    driver_id = re.search(r'Driver ID (\d+)', pptx_path.name)
    if not driver_id:
        print(f'  ⚠ Não consegui extrair Driver ID de {pptx_path.name}')
        return

    did = driver_id.group(1)
    rows = df_driver.get(did, [])
    n = len(rows)

    print(f'\nDriver {did}: {n} SHPs no CSV')

    if n == 0:
        print(f'  → nenhum dado no CSV, pulando')
        return

    prs = Presentation(str(pptx_path))
    total_data_slides = max(1, -(-n // ROWS_PER_SLIDE))  # ceil division

    # Slide 2 (índice 1) é o slide de dados template
    # Atualizar slide 2 com as primeiras 32 linhas
    slide2 = prs.slides[1]
    preencher_slide(slide2, rows[:ROWS_PER_SLIDE], 1, total_data_slides)

    # Adicionar slides extras para os dados restantes
    # Estrutura original: 0=header, 1=data, 2=footer
    # Após add_slide: novos slides vão para índices 3,4,...
    # Queremos: 0=header, 1=data, 2..N=continuação, N+1=footer
    for page in range(1, total_data_slides):
        batch = rows[page * ROWS_PER_SLIDE: (page + 1) * ROWS_PER_SLIDE]
        new_slide = copy_slide(prs, 1)
        preencher_slide(new_slide, batch, page + 1, total_data_slides)
        print(f'  ✓ Slide continuação {page+1}/{total_data_slides} ({len(batch)} IDs)')

    # Mover o slide footer (originalmente índice 2) para o final
    if total_data_slides > 1:
        # Footer está no índice 2; mover para o fim
        move_slide(prs, 2, len(prs.slides) - 1)
        print(f'  ✓ Footer movido para slide {len(prs.slides)} (fim)')

    out_path = out_dir / pptx_path.name
    prs.save(str(out_path))
    print(f'  → Salvo: {out_path.name} ({total_data_slides} slide(s) de dados)')

# ── Main ─────────────────────────────────────────────────────────────────────
if __name__ == '__main__':
    print('=' * 60)
    print('COMPLETAR PPTXs — Analise Driver')
    print('=' * 60)

    PASTA_OUT.mkdir(exist_ok=True)
    df = carregar_csv()

    # Agrupar por driver
    df_driver = {}
    for did, grp in df.groupby('DRIVER_ID'):
        if did and str(did).strip():
            df_driver[str(did).strip()] = grp.to_dict('records')

    print(f'Drivers no CSV: {len(df_driver)}')

    pptxs = sorted(PASTA_PPTX.glob('*.pptx'))
    print(f'PPTXs encontrados: {len(pptxs)}')

    ok = 0
    for p in pptxs:
        driver_id = re.search(r'Driver ID (\d+)', p.name)
        if not driver_id:
            continue
        did = driver_id.group(1)
        n_csv = len(df_driver.get(did, []))
        if n_csv > ROWS_PER_SLIDE:
            print(f'\n>>> {p.name}')
            processar(p, df_driver, PASTA_OUT)
            ok += 1
        else:
            print(f'Driver {did}: {n_csv} SHPs → cabe em 1 slide, copiando sem modificação')
            import shutil
            shutil.copy2(str(p), str(PASTA_OUT / p.name))

    print(f'\n{"="*60}')
    print(f'Concluído. {ok} PPTXs com continuação. Saída: {PASTA_OUT}')
